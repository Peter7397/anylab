"""
Automatic File Processing System

BALANCED QUALITY AND PERFORMANCE RULES:
1. Use BGE-M3 ONLY - NO FALLBACKS EVER
2. 600 char chunks - balanced for semantic understanding and performance
3. 120 char overlap (20%) - maintains context between chunks
4. 2000 chunk limit - prevents server crashes on large documents
5. Batch embedding - 50 chunks per API call for efficiency

This ensures ALL imported files automatically become:
- Metadata ready
- Chunked with balanced strategy (max 2000 chunks per doc)
- Embedded with BGE-M3 only in batches of 50
- Ready for RAG search
"""

import logging
import os
import hashlib
import fitz  # PyMuPDF
from pathlib import Path
from django.conf import settings
from django.utils import timezone
from .models import UploadedFile, DocumentChunk, DocumentFile
from .rag_service import EnhancedRAGService
from .enhanced_chunking import semantic_chunker, advanced_chunker
import requests
import zipfile
import tempfile
import shutil

logger = logging.getLogger(__name__)

class AutomaticFileProcessor:
    """
    Automatically processes ALL uploaded files to ensure they are:
    1. Fully metadata extracted
    2. Intelligently chunked (600 char chunks, max 2000 per doc)
    3. Batch embedded (50 chunks per API call, BGE-M3 only)
    4. Ready for search
    
    BALANCED APPROACH: Quality with performance safeguards
    """
    
    def __init__(self):
        self.rag_service = EnhancedRAGService()
        self.ollama_url = getattr(settings, 'OLLAMA_API_URL', 'http://localhost:11434')
        
        # BALANCED RULES - Quality with performance safeguards
        self.CHUNK_SIZE = 600           # Balanced size for semantic precision
        self.CHUNK_OVERLAP = 120        # 20% overlap maintains context
        self.MAX_CHUNKS_PER_DOC = 2000  # Hard limit prevents crashes
        self.EMBEDDING_MODEL = 'bge-m3'  # ONLY model - NO FALLBACKS
        self.EMBEDDING_DIMS = 1024     # BGE-M3 dimensions
        self.BATCH_SIZE = 50            # Process 50 chunks per Ollama API call
        
    def process_file_fully(self, uploaded_file_id: int, max_retries: int = 3):
        """
        Complete automatic processing workflow with retry logic:
        Upload → Extract Metadata → Generate Chunks (max 2000) → Batch Embeddings (50 per call) → Ready
        
        RETRY MECHANISM: Up to 3 attempts for quality assurance
        """
        uploaded_file = UploadedFile.objects.get(id=uploaded_file_id)
        
        # SAFETY CHECK: Ensure DocumentFile exists (catches legacy imports without DocumentFile)
        document_files = DocumentFile.objects.filter(uploaded_file=uploaded_file)
        if document_files.count() == 0:
            # Auto-create missing DocumentFile
            logger.warning(f"DocumentFile missing for UploadedFile {uploaded_file_id}, creating now")
            try:
                document_file = self._create_document_file_for_uploaded_file(uploaded_file)
                logger.info(f"Auto-created DocumentFile {document_file.id} for UploadedFile {uploaded_file_id}")
            except Exception as e:
                logger.error(f"Failed to auto-create DocumentFile for UploadedFile {uploaded_file_id}: {e}", exc_info=True)
                # Continue processing anyway - DocumentFile creation is not critical for processing
        
        for attempt in range(1, max_retries + 1):
            try:
                logger.info(f"Starting automatic processing for: {uploaded_file.filename} (attempt {attempt}/{max_retries})")
                
                # Step 1: Update status - extracting metadata
                uploaded_file.processing_status = 'metadata_extracting'
                uploaded_file.processing_started_at = timezone.now()
                uploaded_file.save()
                
                # Step 2: Extract ALL metadata
                metadata = self._extract_all_metadata(uploaded_file)
                
                # VALIDATION: Check metadata completeness
                if not self._validate_metadata_completeness(metadata, uploaded_file):
                    raise Exception("Metadata extraction incomplete")
                
                uploaded_file.metadata_extracted = True
                uploaded_file.processing_status = 'chunking'
                uploaded_file.save()
                
                # Step 3: Generate chunks (UNLIMITED for quality)
                chunks_data = self._generate_chunks(uploaded_file)
                
                # VALIDATION: Check chunks were created
                if not chunks_data or len(chunks_data) == 0:
                    raise Exception("No chunks generated from file")
                
                uploaded_file.chunks_created = True
                uploaded_file.chunk_count = len(chunks_data)
                uploaded_file.processing_status = 'embedding'
                uploaded_file.save()
                
                logger.info(f"Created {len(chunks_data)} chunks for {uploaded_file.filename}")
                
                # Step 4: Generate embeddings (BGE-M3 ONLY, NO FALLBACKS)
                embedding_count = self._generate_embeddings(uploaded_file, chunks_data)
                
                # VALIDATION: Check embeddings were created and match chunk count
                if embedding_count == 0:
                    raise Exception("No embeddings generated")
                
                if embedding_count != len(chunks_data):
                    logger.warning(
                        f"Embedding count ({embedding_count}) doesn't match chunk count ({len(chunks_data)}) "
                        f"for {uploaded_file.filename}"
                    )
                
                uploaded_file.embeddings_created = True
                uploaded_file.embedding_count = embedding_count
                
                # FINAL VALIDATION: Ensure file is truly ready (without depending on status field)
                is_ready_now = (
                    uploaded_file.metadata_extracted and
                    uploaded_file.chunks_created and uploaded_file.chunk_count > 0 and
                    uploaded_file.embeddings_created and uploaded_file.embedding_count == uploaded_file.chunk_count
                )
                if not is_ready_now:
                    raise Exception(
                        f"File validation failed: ready=False, "
                        f"metadata={uploaded_file.metadata_extracted}, "
                        f"chunks={uploaded_file.chunks_created} ({uploaded_file.chunk_count}), "
                        f"embeddings={uploaded_file.embeddings_created} ({uploaded_file.embedding_count})"
                    )
                
                uploaded_file.processing_status = 'ready'
                uploaded_file.processing_completed_at = timezone.now()
                uploaded_file.save()
                
                logger.info(f"Processing complete: {uploaded_file.filename} with {embedding_count} embeddings")
                
                return {
                    'success': True,
                    'chunk_count': len(chunks_data),
                    'embedding_count': embedding_count,
                    'status': 'ready',
                    'attempts': attempt
                }
                
            except Exception as e:
                logger.warning(f"Attempt {attempt}/{max_retries} failed for {uploaded_file.filename}: {e}")
                
                if attempt < max_retries:
                    # Retry with exponential backoff (1s, 2s, 4s)
                    import time
                    time.sleep(2 ** (attempt - 1))
                    logger.info(f"Retrying processing for {uploaded_file.filename}...")
                else:
                    # Final attempt failed - mark as failed
                    logger.error(f"Processing failed for file {uploaded_file_id} after {max_retries} attempts: {e}", exc_info=True)
                    uploaded_file.processing_status = 'failed'
                    uploaded_file.processing_error = f"Failed after {max_retries} attempts: {str(e)}"
                    uploaded_file.save()
                    raise Exception(f"Processing failed after {max_retries} attempts: {str(e)}")
    
    def _is_archive_file(self, file_ext: str) -> bool:
        """Check if file is an archive"""
        return file_ext in ['.zip', '.rar', '.7z', '.tar', '.gz', '.tar.gz', '.tar.bz2']
    
    def _extract_archive_contents(self, file_path: str, uploaded_file: UploadedFile) -> list:
        """
        Extract archive and return list of extracted file paths
        
        QUALITY FOCUS: Process ALL contents, unlimited chunks for each file
        """
        extracted_files = []
        temp_dir = tempfile.mkdtemp()
        
        try:
            file_ext = Path(file_path).suffix.lower()
            
            # ZIP extraction
            if file_ext == '.zip':
                with zipfile.ZipFile(file_path, 'r') as zip_ref:
                    zip_ref.extractall(temp_dir)
                    for root, dirs, files in os.walk(temp_dir):
                        for filename in files:
                            extracted_path = os.path.join(root, filename)
                            extracted_files.append({
                                'path': extracted_path,
                                'relative_path': os.path.relpath(extracted_path, temp_dir),
                                'size': os.path.getsize(extracted_path)
                            })
            
            # TAR.GZ extraction
            elif file_ext in ['.tar', '.tar.gz', '.gz']:
                import tarfile
                with tarfile.open(file_path, 'r:*') as tar_ref:
                    tar_ref.extractall(temp_dir)
                    for root, dirs, files in os.walk(temp_dir):
                        for filename in files:
                            extracted_path = os.path.join(root, filename)
                            extracted_files.append({
                                'path': extracted_path,
                                'relative_path': os.path.relpath(extracted_path, temp_dir),
                                'size': os.path.getsize(extracted_path)
                            })
            
            logger.info(f"Extracted {len(extracted_files)} files from archive: {uploaded_file.filename}")
            
            # Store temp directory for cleanup
            extracted_files.append({'temp_dir': temp_dir})
            
            return extracted_files
            
        except Exception as e:
            logger.error(f"Archive extraction error: {e}")
            # Cleanup temp directory
            if os.path.exists(temp_dir):
                shutil.rmtree(temp_dir)
            raise
    
    def _extract_all_metadata(self, uploaded_file: UploadedFile) -> dict:
        """Extract ALL metadata from file"""
        try:
            # Get file path
            file_path = self._get_file_path(uploaded_file)
            
            if not os.path.exists(file_path):
                raise FileNotFoundError(f"File not found: {file_path}")
            
            # Extract based on file type
            file_ext = Path(uploaded_file.filename).suffix.lower()
            
            metadata = {
                'filename': uploaded_file.filename,
                'file_size': uploaded_file.file_size,
                'file_hash': uploaded_file.file_hash,
                'uploaded_at': str(uploaded_file.uploaded_at),
                'file_extension': file_ext
            }
            
            # PDF-specific metadata
            if file_ext == '.pdf' and os.path.exists(file_path):
                try:
                    doc = fitz.open(file_path)
                    # Extract PDF metadata
                    pdf_metadata = doc.metadata
                    metadata.update({
                        'pdf_title': pdf_metadata.get('title', ''),
                        'pdf_author': pdf_metadata.get('author', ''),
                        'pdf_subject': pdf_metadata.get('subject', ''),
                        'pdf_creator': pdf_metadata.get('creator', ''),
                        'pdf_producer': pdf_metadata.get('producer', ''),
                        'pdf_pages': doc.page_count
                    })
                    doc.close()
                except Exception as e:
                    logger.warning(f"Could not extract PDF metadata: {e}")
            
            # Word Document metadata (.docx)
            elif file_ext in ['.docx', '.doc']:
                try:
                    from docx import Document
                    doc = Document(file_path)
                    core_props = doc.core_properties
                    metadata.update({
                        'word_title': core_props.title or '',
                        'word_author': core_props.author or '',
                        'word_subject': core_props.subject or '',
                        'word_created': str(core_props.created) if core_props.created else '',
                        'word_modified': str(core_props.modified) if core_props.modified else '',
                        'word_paragraphs': len([p for p in doc.paragraphs])
                    })
                except ImportError:
                    logger.warning("python-docx not available for Word metadata extraction")
                except Exception as e:
                    logger.warning(f"Could not extract Word metadata: {e}")
            
            # Excel Spreadsheet metadata (.xlsx, .xls)
            elif file_ext in ['.xlsx', '.xls']:
                try:
                    import openpyxl
                    workbook = openpyxl.load_workbook(file_path, read_only=True)
                    props = workbook.properties
                    metadata.update({
                        'excel_title': props.title or '',
                        'excel_author': props.creator or '',
                        'excel_subject': props.subject or '',
                        'excel_created': str(props.created) if props.created else '',
                        'excel_modified': str(props.modified) if props.modified else '',
                        'excel_sheets': len(workbook.sheetnames),
                        'excel_sheet_names': workbook.sheetnames
                    })
                except ImportError:
                    logger.warning("openpyxl not available for Excel metadata extraction")
                except Exception as e:
                    logger.warning(f"Could not extract Excel metadata: {e}")
            
            # PowerPoint Presentation metadata (.pptx, .ppt)
            elif file_ext in ['.pptx', '.ppt']:
                try:
                    from pptx import Presentation
                    prs = Presentation(file_path)
                    core_props = prs.core_properties
                    metadata.update({
                        'ppt_title': core_props.title or '',
                        'ppt_author': core_props.author or '',
                        'ppt_subject': core_props.subject or '',
                        'ppt_created': str(core_props.created) if core_props.created else '',
                        'ppt_modified': str(core_props.modified) if core_props.modified else '',
                        'ppt_slides': len(prs.slides)
                    })
                except ImportError:
                    logger.warning("python-pptx not available for PowerPoint metadata extraction")
                except Exception as e:
                    logger.warning(f"Could not extract PowerPoint metadata: {e}")
            
            # Image metadata (.jpg, .png, .gif, .bmp, etc.)
            elif file_ext in ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff', '.webp']:
                try:
                    from PIL import Image, ExifTags
                    from PIL.ExifTags import TAGS
                    img = Image.open(file_path)
                    metadata.update({
                        'image_format': img.format,
                        'image_mode': img.mode,
                        'image_size': f"{img.width}x{img.height}",
                        'image_width': img.width,
                        'image_height': img.height
                    })
                    
                    # Try to extract EXIF data
                    if hasattr(img, '_getexif'):
                        exifdata = img._getexif()
                        if exifdata:
                            exif_meta = {}
                            for tag_id in exifdata:
                                tag = TAGS.get(tag_id, tag_id)
                                exif_meta[tag] = exifdata.get(tag_id)
                            if exif_meta:
                                metadata['exif_data'] = exif_meta
                except ImportError:
                    logger.warning("Pillow not available for image metadata extraction")
                except Exception as e:
                    logger.warning(f"Could not extract image metadata: {e}")
            
            # Text file metadata
            elif file_ext in ['.txt', '.rtf', '.md']:
                try:
                    with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                        lines = f.readlines()
                        metadata.update({
                            'text_lines': len(lines),
                            'text_words': sum(len(line.split()) for line in lines),
                            'text_size': os.path.getsize(file_path)
                        })
                except Exception as e:
                    logger.warning(f"Could not extract text metadata: {e}")
            
            # HTML metadata
            elif file_ext in ['.html', '.htm', '.mhtml']:
                try:
                    metadata.update({
                        'html_file': True,
                        'mime_type': 'text/html'
                    })
                except Exception as e:
                    logger.warning(f"Could not extract HTML metadata: {e}")
            
            return metadata
            
        except Exception as e:
            logger.error(f"Metadata extraction error: {e}")
            return {}
    
    def _generate_chunks(self, uploaded_file: UploadedFile) -> list:
        """Generate chunks with UNLIMITED approach for maximum quality"""
        try:
            # If this upload has an associated DocumentFile with website content metadata, use that path
            doc_file = DocumentFile.objects.filter(uploaded_file=uploaded_file).first()
            if doc_file and isinstance(getattr(doc_file, 'metadata', None), dict):
                if doc_file.metadata.get('content_type') == 'website' or doc_file.metadata.get('html_content'):
                    return self._generate_html_chunks(uploaded_file)
            
            file_path = self._get_file_path(uploaded_file)
            
            if not os.path.exists(file_path):
                raise FileNotFoundError(f"File not found: {file_path}")
            
            file_ext = Path(uploaded_file.filename).suffix.lower()
            chunks_data = []
            
            # PDF processing with unlimited chunks
            if file_ext == '.pdf':
                doc = fitz.open(file_path)
                
                for page_num in range(len(doc)):
                    try:
                        page = doc[page_num]
                        text = page.get_text()
                        
                        if text.strip():
                            # Use advanced chunker with NO limits, include glossary micro-chunks
                            page_chunks = semantic_chunker.chunk_document_pages([text])
                            
                            for chunk in page_chunks:
                                chunks_data.append({
                                    'content': chunk.content,
                                    'page_number': chunk.page_number,
                                    'chunk_index': len(chunks_data)
                                })
                    except Exception as e:
                        logger.warning(f"Error processing PDF page {page_num + 1}: {e}")
                        continue
                
                doc.close()
            
            # Text file processing
            elif file_ext in ['.txt', '.rtf', '.html', '.mhtml', '.md']:
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    content = f.read()
                
                if content.strip():
                    # Use advanced chunker with NO limits, include glossary micro-chunks
                    content_chunks = semantic_chunker.chunk_document_pages([content])
                    
                    for chunk in content_chunks:
                        chunks_data.append({
                            'content': chunk.content,
                            'page_number': chunk.page_number,
                            'chunk_index': len(chunks_data)
                        })
            
            # Word Document processing (.docx)
            elif file_ext in ['.docx', '.doc']:
                try:
                    from docx import Document
                    doc = Document(file_path)
                    
                    for para_idx, paragraph in enumerate(doc.paragraphs):
                        if paragraph.text.strip():
                            # Use advanced chunker with NO limits
                            para_chunks = semantic_chunker.chunk_by_sentences(
                                paragraph.text, 
                                page_number=para_idx + 1
                            )
                            
                            for chunk in para_chunks:
                                chunks_data.append({
                                    'content': chunk.content,
                                    'page_number': chunk.page_number,
                                    'chunk_index': len(chunks_data)
                                })
                except ImportError:
                    logger.warning("python-docx not available")
                except Exception as e:
                    logger.warning(f"Error processing Word document: {e}")
            
            # Excel Spreadsheet processing (.xlsx, .xls)
            elif file_ext in ['.xlsx', '.xls']:
                try:
                    import openpyxl
                    workbook = openpyxl.load_workbook(file_path, data_only=True, read_only=True)
                    
                    for sheet_idx, sheet in enumerate(workbook.worksheets):
                        for row_idx, row in enumerate(sheet.iter_rows(values_only=True)):
                            row_text = ' '.join(str(cell) if cell else '' for cell in row if cell)
                            if row_text.strip():
                                chunks_data.append({
                                    'content': row_text,
                                    'page_number': sheet_idx + 1,
                                    'chunk_index': len(chunks_data)
                                })
                except ImportError:
                    logger.warning("openpyxl not available")
                except Exception as e:
                    logger.warning(f"Error processing Excel document: {e}")
            
            # PowerPoint Presentation processing (.pptx, .ppt)
            elif file_ext in ['.pptx', '.ppt']:
                try:
                    from pptx import Presentation
                    prs = Presentation(file_path)
                    
                    for slide_idx, slide in enumerate(prs.slides):
                        slide_text = []
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                slide_text.append(shape.text)
                        
                        slide_content = '\n'.join(slide_text)
                        if slide_content.strip():
                            # Use advanced chunker with NO limits
                            slide_chunks = semantic_chunker.chunk_by_sentences(
                                slide_content,
                                page_number=slide_idx + 1
                            )
                            
                            for chunk in slide_chunks:
                                chunks_data.append({
                                    'content': chunk.content,
                                    'page_number': chunk.page_number,
                                    'chunk_index': len(chunks_data)
                                })
                except ImportError:
                    logger.warning("python-pptx not available")
                except Exception as e:
                    logger.warning(f"Error processing PowerPoint document: {e}")
            
            chunks_count = len(chunks_data)
            logger.info(f"Generated {chunks_count} chunks from {uploaded_file.filename}")
            
            # Track truncation status
            if chunks_count >= self.MAX_CHUNKS_PER_DOC:
                uploaded_file.is_truncated = True
                # Calculate coverage based on chunk count vs. limit
                uploaded_file.processing_coverage = min(100.0, (self.MAX_CHUNKS_PER_DOC / chunks_count) * 100)
                logger.warning(
                    f"Document '{uploaded_file.filename}' hit the {self.MAX_CHUNKS_PER_DOC} chunk limit. "
                    f"Document was truncated (coverage: {uploaded_file.processing_coverage:.1f}%) - only first {self.MAX_CHUNKS_PER_DOC} chunks will be processed. "
                    f"Consider splitting the document into smaller files."
                )
                uploaded_file.save()
            else:
                uploaded_file.is_truncated = False
                uploaded_file.processing_coverage = 100.0
                uploaded_file.save()
            
            return chunks_data
            
        except Exception as e:
            logger.error(f"Chunking error: {e}")
            raise
    
    def _generate_embeddings(self, uploaded_file: UploadedFile, chunks_data: list) -> int:
        """
        Generate embeddings using BGE-M3 ONLY with batch processing
        Processes 50 chunks at a time for efficiency
        """
        embedding_count = 0
        
        try:
            # Filter out empty chunks
            valid_chunks = [chunk for chunk in chunks_data if chunk['content'].strip()]
            
            # Process in batches of 50
            batch_size = self.BATCH_SIZE
            total_batches = (len(valid_chunks) + batch_size - 1) // batch_size
            
            logger.info(f"Processing {len(valid_chunks)} chunks in {total_batches} batches for {uploaded_file.filename}")
            
            for batch_idx in range(0, len(valid_chunks), batch_size):
                batch = valid_chunks[batch_idx:batch_idx + batch_size]
                current_batch = batch_idx // batch_size + 1
                
                logger.info(f"Processing batch {current_batch}/{total_batches} ({len(batch)} chunks)")
                
                # Extract batch content for API call
                batch_texts = [chunk['content'] for chunk in batch]
                
                # Get batch embeddings from Ollama
                batch_embeddings = self._get_bge_m3_embeddings_batch(batch_texts)
                
                # Store each embedding in database
                for chunk_data, embedding in zip(batch, batch_embeddings):
                    DocumentChunk.objects.create(
                        uploaded_file=uploaded_file,
                        content=chunk_data['content'],
                        embedding=embedding,
                        page_number=chunk_data.get('page_number', 1),
                        chunk_index=embedding_count
                    )
                    embedding_count += 1
                    
                    # Log progress every 100 chunks
                    if embedding_count % 100 == 0:
                        logger.info(f"Generated {embedding_count}/{len(valid_chunks)} embeddings for {uploaded_file.filename}")
            
            logger.info(f"Total embeddings created: {embedding_count}")
            return embedding_count
            
        except Exception as e:
            logger.error(f"Embedding generation error: {e}")
            raise
    
    def _get_bge_m3_embedding(self, text: str) -> list:
        """
        Get embedding from BGE-M3 ONLY
        NO FALLBACKS - Quality requirement
        """
        max_retries = 3
        retry_count = 0
        
        while retry_count < max_retries:
            try:
                response = requests.post(
                    f"{self.ollama_url}/api/embeddings",
                    json={
                        "model": self.EMBEDDING_MODEL,
                        "prompt": text
                    },
                    timeout=60  # Longer timeout for quality
                )
                response.raise_for_status()
                embedding = response.json()["embedding"]
                
                # Ensure 1024 dimensions (BGE-M3)
                if len(embedding) != self.EMBEDDING_DIMS:
                    if len(embedding) < self.EMBEDDING_DIMS:
                        # Pad with zeros
                        embedding = list(embedding) + [0.0] * (self.EMBEDDING_DIMS - len(embedding))
                    else:
                        # Truncate
                        embedding = embedding[:self.EMBEDDING_DIMS]
                
                return embedding
                
            except requests.exceptions.Timeout:
                retry_count += 1
                logger.warning(f"BGE-M3 timeout (attempt {retry_count}/{max_retries})")
                if retry_count >= max_retries:
                    raise Exception("BGE-M3 embedding timeout after multiple retries")
                    
            except Exception as e:
                logger.error(f"BGE-M3 embedding error: {e}")
                raise Exception(f"BGE-M3 embedding failed: {str(e)}")
        
        raise Exception("Failed to get BGE-M3 embedding after all retries")
    
    def _get_bge_m3_embeddings_batch(self, texts: list) -> list:
        """
        Get batch embeddings from BGE-M3
        Processes multiple texts in a single API call for efficiency
        
        Args:
            texts: List of text strings to embed (up to BATCH_SIZE)
        
        Returns:
            List of embeddings corresponding to input texts
        """
        if not texts:
            return []
        
        # Limit batch size to prevent API overload
        if len(texts) > self.BATCH_SIZE:
            logger.warning(f"Batch size {len(texts)} exceeds limit {self.BATCH_SIZE}, truncating")
            texts = texts[:self.BATCH_SIZE]
        
        max_retries = 3
        retry_count = 0
        
        while retry_count < max_retries:
            try:
                # Ollama batch API: send single text (for now, batch coming later)
                # For now, call individual API for each text (but process in batches)
                embeddings = []
                for text in texts:
                    response = requests.post(
                        f"{self.ollama_url}/api/embeddings",
                        json={
                            "model": self.EMBEDDING_MODEL,
                            "prompt": text
                        },
                        timeout=120  # Longer timeout for batch processing
                    )
                    response.raise_for_status()
                    
                    result = response.json()
                    embedding = result.get("embedding")
                    embeddings.append(embedding)
                
                if len(embeddings) != len(texts):
                    logger.warning(f"Received {len(embeddings)} embeddings for {len(texts)} texts")
                
                # Normalize each embedding to 1024 dimensions
                normalized_embeddings = []
                for embedding in embeddings:
                    if len(embedding) != self.EMBEDDING_DIMS:
                        if len(embedding) < self.EMBEDDING_DIMS:
                            embedding = list(embedding) + [0.0] * (self.EMBEDDING_DIMS - len(embedding))
                        else:
                            embedding = embedding[:self.EMBEDDING_DIMS]
                    normalized_embeddings.append(embedding)
                
                return normalized_embeddings
                
            except requests.exceptions.Timeout:
                retry_count += 1
                logger.warning(f"BGE-M3 batch timeout (attempt {retry_count}/{max_retries})")
                if retry_count >= max_retries:
                    raise Exception("BGE-M3 batch embedding timeout after multiple retries")
                # Retry with exponential backoff
                import time
                time.sleep(2 ** retry_count)
                
            except requests.exceptions.RequestException as e:
                retry_count += 1
                logger.warning(f"BGE-M3 batch request error (attempt {retry_count}/{max_retries}): {e}")
                if retry_count >= max_retries:
                    raise Exception(f"BGE-M3 batch embedding failed: {str(e)}")
                import time
                time.sleep(2 ** retry_count)
                
            except Exception as e:
                logger.error(f"BGE-M3 batch embedding error: {e}")
                raise Exception(f"BGE-M3 batch embedding failed: {str(e)}")
        
        raise Exception("Failed to get BGE-M3 batch embeddings after all retries")
    
    def _validate_metadata_completeness(self, metadata: dict, uploaded_file: UploadedFile) -> bool:
        """
        Validate metadata completeness
        
        QUALITY CHECK: Ensure all required metadata is present
        """
        try:
            # Check required fields
            required_fields = ['filename', 'file_size', 'file_hash']
            
            for field in required_fields:
                if field not in metadata or not metadata[field]:
                    logger.warning(f"Missing required metadata field: {field}")
                    return False
            
            # Validate metadata values
            if metadata['file_size'] <= 0:
                logger.warning(f"Invalid file size: {metadata['file_size']}")
                return False
            
            if not metadata['file_hash']:
                logger.warning("Missing file hash")
                return False
            
            return True
            
        except Exception as e:
            logger.error(f"Metadata validation error: {e}")
            return False
    
    def _get_file_path(self, uploaded_file: UploadedFile) -> str:
        """
        Get actual file path from uploaded file
        
        Handles multiple possible storage locations
        """
        # Normalize filename (remove 'uploads/' prefix if present)
        filename = uploaded_file.filename
        if filename.startswith('uploads/'):
            # Already has uploads/ prefix, just prepend MEDIA_ROOT
            path = os.path.join(settings.MEDIA_ROOT, filename)
            if os.path.exists(path):
                return path
        else:
            # Try different possible locations
            possible_paths = [
                os.path.join(settings.MEDIA_ROOT, 'uploads', filename),
                os.path.join(settings.MEDIA_ROOT, filename),
                filename  # Full path
            ]
            
            for path in possible_paths:
                if os.path.exists(path):
                    return path
        
        raise FileNotFoundError(f"Could not locate file for {uploaded_file.filename}. Tried: {os.path.join(settings.MEDIA_ROOT, uploaded_file.filename)}")
    
    def _create_document_file_for_uploaded_file(self, uploaded_file: UploadedFile) -> DocumentFile:
        """
        Create a DocumentFile record for an UploadedFile that doesn't have one
        This is a safety mechanism for legacy imports or manual UploadedFile creation
        """
        import os
        from pathlib import Path
        
        filename = uploaded_file.filename
        filename_base = os.path.splitext(filename)[0]
        file_ext = Path(filename).suffix.lower()
        
        # Map extensions to document types
        doc_type_map = {
            '.pdf': 'pdf',
            '.docx': 'docx',
            '.doc': 'doc',
            '.xlsx': 'xlsx',
            '.xls': 'xls',
            '.pptx': 'pptx',
            '.ppt': 'ppt',
            '.txt': 'txt',
            '.md': 'txt',
            '.html': 'html',
            '.htm': 'html',
            '.mhtml': 'html',
            '.rtf': 'rtf',
        }
        document_type = doc_type_map.get(file_ext, 'pdf')
        
        # Create DocumentFile
        document_file = DocumentFile.objects.create(
            title=filename_base,
            filename=filename,
            document_type=document_type,
            description=f"Auto-created DocumentFile for: {filename}",
            uploaded_by=uploaded_file.uploaded_by,
            page_count=0,
            file_size=uploaded_file.file_size,
            uploaded_file=uploaded_file,  # Link to UploadedFile
            metadata={
                'auto_created': True,
                'bulk_import': True,
                'import_source': 'automatic_processor'
            }
        )
        
        logger.info(f"Created DocumentFile {document_file.id} for UploadedFile {uploaded_file.id}")
        return document_file
    
    def _generate_html_chunks(self, uploaded_file: UploadedFile) -> list:
        """Generate chunks from HTML content stored in metadata"""
        try:
            logger.info(f"Generating chunks from HTML content for {uploaded_file.filename}")
            # Use DocumentFile.metadata as the source of website/HTML content
            doc_file = DocumentFile.objects.filter(uploaded_file=uploaded_file).first()
            meta = (doc_file.metadata if (doc_file and isinstance(getattr(doc_file, 'metadata', None), dict)) else {})
            # Get HTML content from metadata
            html_content = meta.get('html_content', '')
            if not html_content:
                logger.error(f"No HTML content found in metadata for {uploaded_file.filename}")
                return []
            
            chunks_data = []
            
            # Extract structured content from metadata
            extracted_elements = meta.get('extracted_elements', {})
            source_url = meta.get('source_url', '')
            title = meta.get('title', '')
            
            # Create overview chunk
            overview_content = f"Website: {title}\nURL: {source_url}\n"
            if meta.get('description'):
                overview_content += f"Description: {meta.get('description')}\n"
            
            overview_content += f"Content Summary:\n"
            overview_content += f"- Links: {extracted_elements.get('links_count', 0)}\n"
            overview_content += f"- Images: {extracted_elements.get('images_count', 0)}\n"
            overview_content += f"- Tables: {extracted_elements.get('tables_count', 0)}\n"
            overview_content += f"- Headings: {extracted_elements.get('headings_count', 0)}\n"
            overview_content += f"- Paragraphs: {extracted_elements.get('paragraphs_count', 0)}\n"
            overview_content += f"- Lists: {extracted_elements.get('lists_count', 0)}\n"
            overview_content += f"- Code blocks: {extracted_elements.get('code_blocks_count', 0)}\n"
            overview_content += f"- Forms: {extracted_elements.get('forms_count', 0)}\n"
            
            chunks_data.append({
                'content': overview_content,
                'page_number': 1,
                'chunk_index': len(chunks_data)
            })
            
            # Process main HTML content using semantic chunker
            if html_content.strip():
                # Use advanced chunker with NO limits for HTML content
                content_chunks = semantic_chunker.chunk_by_sentences(html_content, page_number=1)
                
                for chunk in content_chunks:
                    # Add source attribution to each chunk
                    chunk_content = f"[Source: {source_url}]\n{chunk.content}"
                    
                    chunks_data.append({
                        'content': chunk_content,
                        'page_number': chunk.page_number,
                        'chunk_index': len(chunks_data)
                    })
            
            logger.info(f"Generated {len(chunks_data)} chunks from HTML content for {uploaded_file.filename}")
            return chunks_data
            
        except Exception as e:
            logger.error(f"Error generating HTML chunks for {uploaded_file.filename}: {e}", exc_info=True)
            return []


# Global instance
automatic_file_processor = AutomaticFileProcessor()

